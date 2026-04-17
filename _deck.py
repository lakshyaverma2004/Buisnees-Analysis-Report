"""
Big 4 Executive Deck — Explicit Q&A Mapping
Each slide clearly labels which task question it answers
Broader analysis beyond just cancellations
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY     = RGBColor(20, 33, 61)
SLATE    = RGBColor(45, 55, 72)
WHITE    = RGBColor(255, 255, 255)
LIGHT    = RGBColor(200, 205, 215)
MUTED    = RGBColor(130, 140, 155)
ACCENT   = RGBColor(0, 180, 140)
WARN     = RGBColor(230, 126, 86)
GOLD     = RGBColor(218, 185, 107)
DIVIDER  = RGBColor(55, 70, 95)
QTAG     = RGBColor(120, 180, 255)    # blue tag for question labels

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

def line(slide, left, top, width, color=ACCENT):
    s = slide.shapes.add_shape(1, left, top, width, Inches(0.035))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def vline(slide, left, top, height, color=DIVIDER):
    s = slide.shapes.add_shape(1, left, top, Inches(0.02), height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def panel(slide, left, top, width, height, color=SLATE):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def txt(slide, left, top, w, h, text, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    t = slide.shapes.add_textbox(left, top, w, h)
    t.text_frame.word_wrap = True
    p = t.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.alignment = align
    return t

def bullets(slide, left, top, w, h, items):
    t = slide.shapes.add_textbox(left, top, w, h)
    t.text_frame.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(items):
        p = t.text_frame.paragraphs[0] if i == 0 else t.text_frame.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.bold = bld; p.space_after = Pt(4)
    return t

def img(slide, path, left, top, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, w, h)

def slide_header(slide, qtag, title, subtitle=""):
    # Question tag
    panel(slide, Inches(0.7), Inches(0.2), Inches(3.5), Inches(0.32), RGBColor(30, 55, 100))
    txt(slide, Inches(0.8), Inches(0.2), Inches(3.3), Inches(0.32), qtag, 10, QTAG, True)
    # Title
    txt(slide, Inches(0.7), Inches(0.55), Inches(11.8), Inches(0.5), title, 22, WHITE, True)
    line(slide, Inches(0.7), Inches(1.05), Inches(1.2))
    if subtitle:
        txt(slide, Inches(0.7), Inches(1.12), Inches(11.8), Inches(0.35), subtitle, 12, LIGHT)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(1), Inches(2.2), Inches(1.5), ACCENT)
txt(s, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
    "Hotel Bookings Portfolio Review", 42, WHITE, True)
txt(s, Inches(1), Inches(3.6), Inches(11), Inches(0.6),
    "Key Observations  |  Root Cause Analysis  |  Business Recommendations", 20, LIGHT)
line(s, Inches(1), Inches(4.4), Inches(4), DIVIDER)
txt(s, Inches(1), Inches(4.7), Inches(11), Inches(0.4),
    "CONFIDENTIAL", 11, MUTED, True)
txt(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
    "Analysis Period: April 2024 \u2013 April 2025  |  30,000 Transactions  |  3 Channels \u00b7 3 Room Types \u00b7 15 Cities", 12, MUTED)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Q1.1: THREE MEANINGFUL TRENDS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 1.1  |  Highlight at least 3 meaningful trends or patterns in the data",
    "The Platform Is Growing Fast, But Operational Leaks Are Scaling With It",
    "Three structural patterns define the current portfolio: revenue momentum, a concentrated cancellation problem, and untapped customer loyalty.")

# Trend 1
panel(s, Inches(0.7), Inches(1.5), Inches(3.8), Inches(5.5), SLATE)
txt(s, Inches(0.9), Inches(1.6), Inches(3.4), Inches(0.3),
    "TREND 1: REVENUE IS ACCELERATING", 10, ACCENT, True)
txt(s, Inches(0.9), Inches(1.95), Inches(1.8), Inches(0.5),
    "+23%", 36, GOLD, True)
txt(s, Inches(2.7), Inches(2.05), Inches(1.5), Inches(0.3),
    "YoY growth in avg\nbooking value", 10, MUTED)
bullets(s, Inches(0.9), Inches(2.7), Inches(3.4), Inches(4.0), [
    ("Avg booking value grew from $22.5K (Apr 2024)", 11, LIGHT, False),
    ("to $27.8K (Apr 2025).", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Total confirmed revenue: $638.8M across 21,672", 11, LIGHT, False),
    ("confirmed bookings.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("The platform is moving upmarket \u2014 either attracting", 11, LIGHT, False),
    ("premium travelers or implementing better pricing.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Monthly volume is stable (~2,200\u20132,500/month)", 11, WHITE, True),
    ("while value per booking climbs consistently.", 11, LIGHT, False),
])

# Trend 2
panel(s, Inches(4.7), Inches(1.5), Inches(3.8), Inches(5.5), SLATE)
txt(s, Inches(4.9), Inches(1.6), Inches(3.4), Inches(0.3),
    "TREND 2: CANCELLATIONS ARE CONCENTRATED", 10, ACCENT, True)
txt(s, Inches(4.9), Inches(1.95), Inches(1.8), Inches(0.5),
    "83%", 36, WARN, True)
txt(s, Inches(6.5), Inches(2.05), Inches(1.5), Inches(0.3),
    "of cancellations\nlack check-in dates", 10, MUTED)
bullets(s, Inches(4.9), Inches(2.7), Inches(3.4), Inches(4.0), [
    ("20.2% overall cancel rate (6,070 of 30,000).", 11, LIGHT, False),
    ("But this is NOT evenly spread.", 11, WHITE, True),
    ("", 4, LIGHT, False),
    ("5,468 bookings had no check-in date assigned.", 11, LIGHT, False),
    ("92.3% of those were cancelled.", 11, WARN, True),
    ("", 4, LIGHT, False),
    ("Travel Agents: 27.9% cancel rate.", 11, LIGHT, False),
    ("Standard rooms: 23.3% cancel rate.", 11, LIGHT, False),
    ("Summer months: up to 30.3% cancel rate.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("The problem is specific, not systemic.", 11, WHITE, True),
])

# Trend 3
panel(s, Inches(8.7), Inches(1.5), Inches(4.0), Inches(5.5), SLATE)
txt(s, Inches(8.9), Inches(1.6), Inches(3.6), Inches(0.3),
    "TREND 3: HIGH CUSTOMER CONCENTRATION", 10, ACCENT, True)
txt(s, Inches(8.9), Inches(1.95), Inches(1.8), Inches(0.5),
    "60x", 36, GOLD, True)
txt(s, Inches(10.5), Inches(2.05), Inches(1.8), Inches(0.3),
    "avg bookings per\nunique customer", 10, MUTED)
bullets(s, Inches(8.9), Inches(2.7), Inches(3.6), Inches(4.0), [
    ("Only 499 unique customers across 30,000 bookings.", 11, LIGHT, False),
    ("That\u2019s ~60 bookings per customer on average.", 11, WHITE, True),
    ("", 4, LIGHT, False),
    ("This is a repeat-heavy business. These are not", 11, LIGHT, False),
    ("one-time leisure travelers \u2014 likely corporate", 11, LIGHT, False),
    ("or travel management accounts.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Opportunity: loyalty tiers, volume discounts,", 11, ACCENT, True),
    ("and personalized pricing are highly viable.", 11, ACCENT, False),
    ("", 4, LIGHT, False),
    ("Risk: losing even 1 heavy customer = losing", 11, WARN, False),
    ("~$1.5M in annual bookings.", 11, WARN, True),
])

txt(s, Inches(0.7), Inches(7.15), Inches(11.9), Inches(0.2),
    "Source: Hotel_bookings_final.csv | 30,000 transactions | Apr 2024\u2013Apr 2025", 8, MUTED)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Q1.2: BOOKING PATTERNS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 1.2  |  How do booking patterns vary across channels, room types, and star ratings?",
    "Web Drives Value, Standard Rooms Drive Volume, Star Ratings Are a Non-Factor",
    "Booking behavior is shaped primarily by channel and room type \u2014 star rating has minimal influence on any metric.")

# Channel column
panel(s, Inches(0.7), Inches(1.5), Inches(3.8), Inches(3.5), SLATE)
txt(s, Inches(0.9), Inches(1.55), Inches(3.4), Inches(0.25), "BY CHANNEL", 10, ACCENT, True)
bullets(s, Inches(0.9), Inches(1.85), Inches(3.4), Inches(2.8), [
    ("Web: 50% of bookings | $28,191 avg value", 11, WHITE, True),
    ("Highest value, lowest cancel. Self-selecting users.", 10, LIGHT, False),
    ("17.6% cancel rate. The golden channel.", 10, ACCENT, False),
    ("", 3, LIGHT, False),
    ("Mobile App: 40% of bookings | $21,351 avg", 11, WHITE, True),
    ("24% less value than Web. Impulse-driven.", 10, LIGHT, False),
    ("21.6% cancel rate. Under-monetized, not broken.", 10, LIGHT, False),
    ("", 3, LIGHT, False),
    ("Travel Agent: 10% of bookings | $24,454 avg", 11, WHITE, True),
    ("Mid-value but 27.9% cancel. Speculative booking.", 10, WARN, False),
])

# Room column
panel(s, Inches(4.7), Inches(1.5), Inches(3.8), Inches(3.5), SLATE)
txt(s, Inches(4.9), Inches(1.55), Inches(3.4), Inches(0.25), "BY ROOM TYPE", 10, ACCENT, True)
bullets(s, Inches(4.9), Inches(1.85), Inches(3.4), Inches(2.8), [
    ("Standard: 55.2% of bookings | 23.3% cancel", 11, WHITE, True),
    ("The default \u201cplaceholder\u201d choice. Low commitment.", 10, LIGHT, False),
    ("Highest volume but worst retention.", 10, WARN, False),
    ("", 3, LIGHT, False),
    ("Deluxe: 34.9% of bookings | 16.0% cancel", 11, WHITE, True),
    ("Similar price to Standard but 7.3pp lower cancel.", 10, LIGHT, False),
    ("People who choose Deluxe are committed travelers.", 10, ACCENT, False),
    ("", 3, LIGHT, False),
    ("Suite: 9.9% of bookings | 18.0% cancel", 11, WHITE, True),
    ("Premium niche. Stable and predictable behavior.", 10, LIGHT, False),
])

# Star column
panel(s, Inches(8.7), Inches(1.5), Inches(4.0), Inches(3.5), SLATE)
txt(s, Inches(8.9), Inches(1.55), Inches(3.6), Inches(0.25), "BY STAR RATING", 10, ACCENT, True)
bullets(s, Inches(8.9), Inches(1.85), Inches(3.6), Inches(2.8), [
    ("Cancel rates are essentially flat across stars:", 11, WHITE, True),
    ("", 3, LIGHT, False),
    ("\u2022  2-star: 19.8%    \u2022  3-star: 20.2%", 11, LIGHT, False),
    ("\u2022  4-star: 20.0%    \u2022  5-star: 21.3%", 11, LIGHT, False),
    ("", 3, LIGHT, False),
    ("Booking values do scale with stars:", 11, WHITE, True),
    ("2-star avg ~$18K vs 5-star avg ~$35K.", 10, LIGHT, False),
    ("", 3, LIGHT, False),
    ("Key insight: Cancellation is NOT driven", 11, GOLD, True),
    ("by property quality. It\u2019s driven by channel", 11, GOLD, False),
    ("and commitment level at time of booking.", 11, GOLD, False),
])

img(s, "chart02_cancel_rate_by_channel.png", Inches(0.7), Inches(5.2), Inches(4.0), Inches(2.1))
img(s, "chart09_cancel_by_room.png", Inches(4.8), Inches(5.2), Inches(4.0), Inches(2.1))
img(s, "chart05_cancel_by_star.png", Inches(8.9), Inches(5.2), Inches(3.8), Inches(2.1))

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Q1.3: CANCELLATION BEHAVIOR
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 1.3  |  What do you notice about cancellation behavior?",
    "Cancellations Are Predictable: They Cluster Around Missing Data, Agent Bookings, and Peak Months",
    "Five distinct behavioral patterns emerge from the cancellation data.")

panel(s, Inches(0.7), Inches(1.5), Inches(6.2), Inches(5.5), SLATE)
txt(s, Inches(0.9), Inches(1.6), Inches(5.8), Inches(0.25), "FIVE CANCELLATION BEHAVIORS", 10, ACCENT, True)
bullets(s, Inches(0.9), Inches(1.9), Inches(5.8), Inches(5.0), [
    ("\u2776  The \u201cGhost Booking\u201d \u2014 No check-in date assigned", 13, WHITE, True),
    ("5,468 bookings (18.2%) had no check-in date. 92.3% cancelled.", 11, LIGHT, False),
    ("These are effectively dead-on-arrival. 83% of ALL cancellations.", 11, WARN, False),
    ("", 4, LIGHT, False),
    ("\u2777  The \u201cAgent Dump\u201d \u2014 Speculative bulk booking", 13, WHITE, True),
    ("Travel Agents cancel 27.9% of bookings. No deposit, no penalty.", 11, LIGHT, False),
    ("They block rooms, sell what they can, release the rest.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("\u2778  The \u201cPlaceholder\u201d \u2014 Standard room as a backup", 13, WHITE, True),
    ("Standard rooms cancel at 23.3% vs 16.0% for Deluxe.", 11, LIGHT, False),
    ("The cheapest option is the easiest to walk away from.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("\u2779  The \u201cSummer Flipper\u201d \u2014 Peak-season over-booking", 13, WHITE, True),
    ("Jul-Aug cancellation spikes to 28.8\u201330.3%.", 11, LIGHT, False),
    ("Travelers book 2\u20133 hotels, then pick the cheapest closer to the date.", 11, LIGHT, False),
    ("", 4, LIGHT, False),
    ("\u277a  The \u201cCoupon Neutral\u201d \u2014 Promotions don\u2019t cause cancels", 13, WHITE, True),
    ("Coupon users and non-users cancel at nearly identical rates.", 11, LIGHT, False),
    ("Good news: promotions aren\u2019t driving adverse behavior.", 11, ACCENT, False),
])

img(s, "chart17_missing_checkin.png", Inches(7.2), Inches(1.5), Inches(5.5), Inches(2.5))
img(s, "chart20_cancel_by_payment.png", Inches(7.2), Inches(4.2), Inches(5.5), Inches(2.5))

txt(s, Inches(0.7), Inches(7.15), Inches(11.9), Inches(0.2),
    "Source: Hotel_bookings_final.csv | Behavioral cohort analysis", 8, MUTED)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Q2.1: REASONS FOR CANCELLATION PATTERNS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 2.1  |  What might be the reasons for the observed cancellation patterns?",
    "Four Root Causes Explain 95% of Cancellation Behavior",
    "Each root cause maps to a specific, addressable operational gap.")

for i, (cause, why, evidence, fix, clr) in enumerate([
    ("PROCESS GAP:\nMissing Check-in Dates",
     "The booking system allows reservations without assigning a check-in date. These \u201cincomplete\u201d bookings are never operationally activated.",
     "5,468 bookings (18.2%) lack dates.\n92.3% of those cancel.\n83% of all cancellations.",
     "Auto-assign dates at booking.\nForce 48hr confirmation.",
     WARN),
    ("INCENTIVE MISALIGNMENT:\nAgent Booking Model",
     "Travel Agents face zero financial penalty for cancelling. They pre-book rooms speculatively to guarantee supply, then release unsold inventory.",
     "Agent cancel rate: 27.9%.\nWeb cancel rate: 17.6%.\n60% higher cancellation.",
     "Require 20\u201330% non-refundable\ndeposit on agent bookings.",
     WARN),
    ("BEHAVIORAL:\nLow-Commitment Room Choice",
     "Standard is the cheapest room type and acts as a \u201cjust-in-case\u201d placeholder. Guests book it with low intent and cancel when plans firm up.",
     "Standard cancel: 23.3%.\nDeluxe cancel: 16.0%.\n7.3pp gap at similar prices.",
     "Upsell Standard \u2192 Deluxe.\nNon-refundable Standard tiers.",
     GOLD),
    ("MARKET-DRIVEN:\nSeasonal Over-Booking",
     "Peak summer demand (Jul-Aug) enables guests to book multiple hotels and choose last minute. This creates artificial inflation in cancel rates.",
     "Jul: 28.8% | Aug: 30.3%.\nNov (off-peak): 15.9%.\n1.9x seasonal swing.",
     "Summer dynamic pricing +10\u201315%.\nStrict cancel terms in peak.",
     GOLD),
]):
    x = Inches(0.7) + Inches(i * 3.15)
    panel(s, x, Inches(1.5), Inches(2.95), Inches(5.5), SLATE)
    txt(s, x + Inches(0.12), Inches(1.55), Inches(2.7), Inches(0.6), cause, 10, clr, True)
    line(s, x + Inches(0.12), Inches(2.2), Inches(2.7), clr)
    txt(s, x + Inches(0.12), Inches(2.3), Inches(2.7), Inches(0.3), "WHY", 9, MUTED, True)
    txt(s, x + Inches(0.12), Inches(2.55), Inches(2.7), Inches(1.3), why, 10, LIGHT)
    txt(s, x + Inches(0.12), Inches(3.9), Inches(2.7), Inches(0.3), "EVIDENCE", 9, MUTED, True)
    txt(s, x + Inches(0.12), Inches(4.15), Inches(2.7), Inches(1.0), evidence, 10, WHITE, True)
    txt(s, x + Inches(0.12), Inches(5.2), Inches(2.7), Inches(0.3), "FIX", 9, MUTED, True)
    txt(s, x + Inches(0.12), Inches(5.45), Inches(2.7), Inches(0.8), fix, 10, ACCENT, True)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Q2.2: WHY SOME CHANNELS/PROPERTIES PERFORM BETTER
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 2.2  |  Why do some booking channels or property types perform better than others?",
    "Performance Gaps Are Driven by User Intent, Financial Accountability, and Friction",
    "The best-performing segments share a common trait: higher user commitment at the time of booking.")

# Left: Channel performance
panel(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.4), SLATE)
txt(s, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.25), "CHANNEL PERFORMANCE DRIVERS", 10, ACCENT, True)
bullets(s, Inches(0.9), Inches(1.9), Inches(5.4), Inches(4.8), [
    ("Web outperforms because users self-select.", 13, WHITE, True),
    ("Web users research across multiple sites before booking.", 10, LIGHT, False),
    ("By the time they book, they\u2019ve already decided. High intent.", 10, LIGHT, False),
    ("Result: $28K avg value, 17.6% cancel rate.", 10, ACCENT, False),
    ("", 4, LIGHT, False),
    ("Mobile App underperforms on value, not retention.", 13, WHITE, True),
    ("Smaller screens = quicker decisions = cheaper bookings.", 10, LIGHT, False),
    ("$21K avg (24% below Web). But cancel rate (21.6%) is closer.", 10, LIGHT, False),
    ("The App isn\u2019t broken \u2014 it\u2019s under-monetized.", 10, GOLD, True),
    ("", 4, LIGHT, False),
    ("Travel Agents fail because there\u2019s no cost to failing.", 13, WHITE, True),
    ("Zero deposit + free cancellation = zero accountability.", 10, LIGHT, False),
    ("Agents book speculatively, sell what they can, dump the rest.", 10, LIGHT, False),
    ("27.9% cancel rate is a structural incentive problem.", 10, WARN, True),
    ("", 4, LIGHT, False),
    ("Deluxe rooms outperform Standard \u2014 same reason.", 13, WHITE, True),
    ("Choosing Deluxe is an active decision, not a default.", 10, LIGHT, False),
    ("It signals higher travel commitment. 16% vs 23.3% cancel.", 10, LIGHT, False),
])

# Right: Charts
img(s, "chart03_avg_value_by_channel.png", Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5))
img(s, "chart21_heatmap_value.png", Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.5))

txt(s, Inches(0.7), Inches(7.15), Inches(11.9), Inches(0.2),
    "Source: Hotel_bookings_final.csv | Channel & room type cohort analysis", 8, MUTED)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Q2.3: SEASONAL/TEMPORAL TRENDS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 2.3  |  Are there seasonal or temporal trends influencing booking values or stay lengths?",
    "Strong Seasonality Exists in Both Cancellation Rates and Booking Values, With a Clear YoY Growth Trend",
    "The platform shows healthy growth momentum, but seasonal patterns create exploitable pricing windows.")

panel(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(2.5), SLATE)
txt(s, Inches(0.9), Inches(1.6), Inches(5.4), Inches(0.2), "BOOKING VALUE TRENDS", 10, ACCENT, True)
bullets(s, Inches(0.9), Inches(1.85), Inches(5.4), Inches(2.0), [
    ("Avg booking value: $22.5K (Apr 2024) \u2192 $27.8K (Apr 2025)", 12, GOLD, True),
    ("23% YoY growth in revenue per booking.", 11, LIGHT, False),
    ("Monthly total revenue: $50M\u2013$75M, trending upward.", 11, LIGHT, False),
    ("Summer months show highest values ($26\u201328K avg).", 11, LIGHT, False),
    ("Off-peak (Feb-Mar) shows lower values ($24\u201325K avg).", 11, LIGHT, False),
])

panel(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), SLATE)
txt(s, Inches(7.0), Inches(1.6), Inches(5.4), Inches(0.2), "STAY LENGTH & LEAD TIME", 10, ACCENT, True)
bullets(s, Inches(7.0), Inches(1.85), Inches(5.4), Inches(2.0), [
    ("Average stay length: 4.0 days (consistent across months).", 12, WHITE, True),
    ("Average lead time: 30.4 days (booking to check-in).", 11, LIGHT, False),
    ("Cancelled bookings have slightly higher lead time.", 11, LIGHT, False),
    ("Longer lead time = more time to change plans = higher cancel.", 11, WARN, False),
    ("Business vs leisure stays show similar length distributions.", 11, LIGHT, False),
])

img(s, "chart13_monthly_total_value.png", Inches(0.7), Inches(4.2), Inches(6.0), Inches(2.9))
img(s, "chart14_monthly_avg_value.png", Inches(6.8), Inches(4.2), Inches(6.0), Inches(2.9))

txt(s, Inches(0.7), Inches(7.15), Inches(11.9), Inches(0.2),
    "Source: Hotel_bookings_final.csv | Monthly cohort analysis", 8, MUTED)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Q3.1: STRATEGIES TO REDUCE CANCELLATIONS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 3.1  |  Suggest strategies to reduce cancellations",
    "Five Interventions Ordered by Impact Can Reduce Cancellation Rate from 20.2% to <15%")

interventions = [
    ("#1", "AUTO-ASSIGN CHECK-IN DATES", "HIGH",
     "Eliminate the #1 cancellation predictor",
     "83% of cancels lack check-in dates. Auto-assign at booking, force 48hr confirmation, flag unconfirmed.",
     "~$20M"),
    ("#2", "TRAVEL AGENT DEPOSIT POLICY", "HIGH",
     "Create financial accountability for agents",
     "Require 20\u201330% non-refundable deposit. Agents cancel 28% because it costs them nothing today.",
     "~$8M"),
    ("#3", "NON-REFUNDABLE RATE TIERS", "MEDIUM",
     "Lock commitment on Standard rooms",
     "8\u201312% discount for non-refundable Standard bookings. Filters out placeholder bookers.",
     "~$5M"),
    ("#4", "PROGRESSIVE CANCELLATION FEES", "MEDIUM",
     "Discourage late cancellations everywhere",
     "Free <24hrs \u2192 10% (7 days) \u2192 25% (3 days) \u2192 50% (same-day). Industry-standard approach.",
     "~$3M"),
    ("#5", "STANDARD \u2192 DELUXE UPSELL CAMPAIGN", "MEDIUM",
     "Shift volume toward committed room type",
     "$50 upgrade credit. Deluxe cancels 16% vs Standard's 23%. 7.3pp improvement per conversion.",
     "~$1M"),
]

for i, (num, title, priority, rationale, detail, impact) in enumerate(interventions):
    y = Inches(1.5) + Inches(i * 1.05)
    pclr = ACCENT if priority == "HIGH" else GOLD
    panel(s, Inches(0.7), y, Inches(11.9), Inches(0.9), SLATE)
    txt(s, Inches(0.85), y + Inches(0.08), Inches(0.5), Inches(0.3), num, 13, pclr, True)
    txt(s, Inches(1.4), y + Inches(0.05), Inches(3.0), Inches(0.3), title, 11, WHITE, True)
    txt(s, Inches(1.4), y + Inches(0.35), Inches(3.0), Inches(0.3), rationale, 9, LIGHT)
    vline(s, Inches(4.5), y + Inches(0.08), Inches(0.7), DIVIDER)
    txt(s, Inches(4.7), y + Inches(0.08), Inches(6.5), Inches(0.7), detail, 10, LIGHT)
    vline(s, Inches(11.2), y + Inches(0.08), Inches(0.7), DIVIDER)
    txt(s, Inches(11.3), y + Inches(0.1), Inches(1.3), Inches(0.3), impact, 16, GOLD, True, PP_ALIGN.CENTER)
    txt(s, Inches(11.3), y + Inches(0.45), Inches(1.3), Inches(0.3), "est. annual", 8, MUTED, False, PP_ALIGN.CENTER)

panel(s, Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.4), ACCENT)
txt(s, Inches(1.0), Inches(6.82), Inches(5), Inches(0.35),
    "TOTAL PROJECTED ANNUAL RECOVERY", 11, NAVY, True)
txt(s, Inches(10.0), Inches(6.78), Inches(2.5), Inches(0.38),
    "~$37M", 20, NAVY, True, PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Q3.2: PROFITABILITY & REPEAT BOOKINGS
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 3.2  |  Recommend ways to improve profitability and increase repeat bookings",
    "Leverage the 60x Repeat Customer Base and Web Channel Dominance to Build a Loyalty Moat")

strategies = [
    ("TIERED LOYALTY PROGRAMME",
     "499 customers make 60 bookings each on average. This is a high-frequency, repeat business.",
     "\u2022  Bronze (10+ bookings): 5% discount on base rate\n\u2022  Silver (25+ bookings): 10% discount + priority upgrades\n\u2022  Gold (50+ bookings): 15% discount + free late checkout\n\u2022  Personalized pricing based on booking history",
     ACCENT),
    ("WEB-FIRST REVENUE MAXIMIZATION",
     "Web is 50% of volume with 32% higher values than App and the lowest cancel rate.",
     "\u2022  Invest in SEO and direct-booking incentives\n\u2022  \u201cBook Direct\u201d price guarantee vs Agent channels\n\u2022  Web-exclusive bundles: room + breakfast + late checkout\n\u2022  Referral programme for existing Web customers",
     ACCENT),
    ("RE-BOOKING INSTEAD OF REFUNDING",
     "6,070 cancellations/year. Even converting 10% = 607 saved bookings worth ~$15M.",
     "\u2022  Offer 5% discount on alternative dates instead of refund\n\u2022  \u201cChange don\u2019t Cancel\u201d flow in App and Web\n\u2022  Flexible date credit valid for 6 months\n\u2022  Partner hotel swap options within same city",
     GOLD),
    ("MOBILE APP MONETIZATION",
     "App users spend $7K less per booking. Close the gap with value-adds, not discounts.",
     "\u2022  App-exclusive room upgrade offers at checkout\n\u2022  Push notifications for last-minute deals (fills inventory)\n\u2022  In-app upsell: breakfast, spa, airport transfer bundles\n\u2022  Mobile-only loyalty point multiplier (2x points)",
     GOLD),
]

for i, (title, context, actions, clr) in enumerate(strategies):
    col = i % 2
    row = i // 2
    x = Inches(0.7) + Inches(col * 6.2)
    y = Inches(1.5) + Inches(row * 2.8)
    panel(s, x, y, Inches(5.9), Inches(2.6), SLATE)
    txt(s, x + Inches(0.15), y + Inches(0.08), Inches(5.6), Inches(0.3), title, 11, clr, True)
    line(s, x + Inches(0.15), y + Inches(0.38), Inches(5.6), clr)
    txt(s, x + Inches(0.15), y + Inches(0.45), Inches(5.6), Inches(0.55), context, 10, LIGHT)
    txt(s, x + Inches(0.15), y + Inches(1.05), Inches(5.6), Inches(1.4), actions, 10, WHITE)

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Q3.3: PRICING, PROMOTIONS, CHANNEL STRATEGY
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
slide_header(s,
    "TASK 3.3  |  Identify opportunities to optimize pricing, promotions, or channel strategy",
    "Three Optimization Levers: Seasonal Pricing Windows, Safe Promotions, and Channel Rebalancing")

# Pricing
panel(s, Inches(0.7), Inches(1.5), Inches(3.7), Inches(5.3), SLATE)
txt(s, Inches(0.9), Inches(1.6), Inches(3.3), Inches(0.25), "PRICING OPTIMIZATION", 10, ACCENT, True)
line(s, Inches(0.9), Inches(1.88), Inches(3.3), ACCENT)
bullets(s, Inches(0.9), Inches(2.0), Inches(3.3), Inches(4.5), [
    ("Summer Dynamic Pricing (Jul-Aug)", 12, WHITE, True),
    ("Raise rates 10\u201315%. Demand absorbs it.", 10, LIGHT, False),
    ("Cancel rate is 30% anyway \u2014 capture value.", 10, LIGHT, False),
    ("Pair with strict same-week cancel terms.", 10, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Off-Peak Value Windows (Nov, Feb-Apr)", 12, WHITE, True),
    ("Nov: 15.9% cancel \u2014 best efficiency month.", 10, LIGHT, False),
    ("Feb-Apr: 16\u201318% cancel with rising values.", 10, LIGHT, False),
    ("Invest marketing budget in these months.", 10, ACCENT, True),
    ("", 4, LIGHT, False),
    ("Star-Based Pricing Tiers", 12, WHITE, True),
    ("5-star avg value is 2x of 2-star.", 10, LIGHT, False),
    ("But cancel rates are flat. Price confidently.", 10, LIGHT, False),
    ("Premium properties can absorb price increases.", 10, LIGHT, False),
])

# Promotions
panel(s, Inches(4.6), Inches(1.5), Inches(3.7), Inches(5.3), SLATE)
txt(s, Inches(4.8), Inches(1.6), Inches(3.3), Inches(0.25), "PROMOTIONS STRATEGY", 10, ACCENT, True)
line(s, Inches(4.8), Inches(1.88), Inches(3.3), ACCENT)
bullets(s, Inches(4.8), Inches(2.0), Inches(3.3), Inches(4.5), [
    ("Coupons: Safe and Effective", 12, WHITE, True),
    ("Coupon users cancel at the SAME rate", 10, LIGHT, False),
    ("as non-users. No adverse selection.", 10, LIGHT, False),
    ("Use coupons freely for acquisition.", 10, ACCENT, True),
    ("", 4, LIGHT, False),
    ("Targeted Upgrade Vouchers", 12, WHITE, True),
    ("Send Standard-room bookers a Deluxe", 10, LIGHT, False),
    ("upgrade voucher ($50 value).", 10, LIGHT, False),
    ("Shifts them to 16% cancel tier from 23%.", 10, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Cashback Restructuring", 12, WHITE, True),
    ("Cashback should reward completion,", 10, LIGHT, False),
    ("not booking. Release cashback AFTER stay.", 10, LIGHT, False),
    ("Creates financial incentive to follow through.", 10, GOLD, True),
    ("", 4, LIGHT, False),
    ("Re-booking Credits > Refunds", 12, WHITE, True),
    ("5% discount on new dates vs full refund.", 10, LIGHT, False),
])

# Channel
panel(s, Inches(8.5), Inches(1.5), Inches(4.1), Inches(5.3), SLATE)
txt(s, Inches(8.7), Inches(1.6), Inches(3.7), Inches(0.25), "CHANNEL STRATEGY", 10, ACCENT, True)
line(s, Inches(8.7), Inches(1.88), Inches(3.7), ACCENT)
bullets(s, Inches(8.7), Inches(2.0), Inches(3.7), Inches(4.5), [
    ("Web: Protect and Grow", 12, WHITE, True),
    ("Best economics. Invest in direct booking.", 10, LIGHT, False),
    ("SEO, SEM, \u201cPrice Match\u201d guarantee.", 10, LIGHT, False),
    ("Target: grow from 50% to 60% of volume.", 10, ACCENT, True),
    ("", 4, LIGHT, False),
    ("Mobile App: Monetize, Don\u2019t Discount", 12, WHITE, True),
    ("Close $7K value gap with upsells, not cuts.", 10, LIGHT, False),
    ("App-exclusive bundles and flash upgrades.", 10, LIGHT, False),
    ("2x loyalty points for App bookings.", 10, LIGHT, False),
    ("", 4, LIGHT, False),
    ("Travel Agent: Reform or Shrink", 12, WHITE, True),
    ("Non-negotiable: 20\u201330% deposits.", 10, LIGHT, False),
    ("Shorter booking windows (7-day max hold).", 10, LIGHT, False),
    ("Performance-based commission structure.", 10, LIGHT, False),
    ("Shift agent volume to Web/App over time.", 10, WARN, True),
])

img(s, "chart23_coupon_cancel.png", Inches(0.7), Inches(6.95), Inches(3.0), Inches(0.4))
img(s, "chart22_heatmap_cancel.png", Inches(4.0), Inches(6.95), Inches(4.5), Inches(0.4))
img(s, "chart26_avg_value_by_city.png", Inches(8.7), Inches(6.95), Inches(3.9), Inches(0.4))

# ═════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BOTTOM LINE
# ═════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
line(s, Inches(1), Inches(1.3), Inches(1.5), ACCENT)
txt(s, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
    "The Bottom Line", 40, WHITE, True)

bullets(s, Inches(1), Inches(2.5), Inches(11), Inches(4.0), [
    ("This is a growing, profitable platform with a concentrated operational leak.", 17, LIGHT, False),
    ("", 6, LIGHT, False),
    ("The cancellation problem is not about demand, quality, or pricing.", 17, WHITE, True),
    ("It\u2019s about process gaps, misaligned incentives, and missing commitment signals.", 17, WHITE, False),
    ("", 8, LIGHT, False),
    ("Three moves capture 80% of the recovery value:", 15, LIGHT, False),
    ("  \u2022  Auto-assign check-in dates  (\u2248$20M)", 14, ACCENT, True),
    ("  \u2022  Agent deposit requirements  (\u2248$8M)", 14, ACCENT, True),
    ("  \u2022  Non-refundable Standard tiers  (\u2248$5M)", 14, ACCENT, True),
    ("", 6, LIGHT, False),
    ("Beyond cancellations: the 499-customer, 60x-repeat base is a loyalty goldmine.", 15, GOLD, True),
    ("Build tiers. Reward Web. Monetize Mobile. Reform Agents.", 15, GOLD, False),
    ("", 8, LIGHT, False),
    ("Projected impact: ~$37M recaptured  |  Cancel rate: 20% \u2192 <15%  |  Timeline: 90 days", 16, WHITE, True),
])

txt(s, Inches(1), Inches(7.0), Inches(11.3), Inches(0.3),
    "Hotel Bookings Portfolio Review  |  Confidential  |  Prepared for Senior Leadership", 9, MUTED, False, PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════
prs.save('Hotel_Bookings_QA_Deck.pptx')
print(f"DONE: Hotel_Bookings_Executive_Deck.pptx ({len(prs.slides)} slides)")
